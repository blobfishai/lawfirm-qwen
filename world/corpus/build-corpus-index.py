#!/usr/bin/env python3
"""Ingest the Calderwood & Harkness corpus into a filesystem-backed store.

WHY THIS EXISTS. Our world stores every document body inline in
world-v13.json. That works at 352 documents / ~137k tokens. C&H is 9,288 files
and ~108M tokens — roughly 790x our corpus — and inlining it would produce an
unloadable world file. So documents live on disk, text is extracted once, and
the world holds only an index. Tools read by path.

This is the change that makes retrieval-at-scale expressible at all. At 137k
tokens an agent can read our entire world in one context window, so "knowing
when to stop searching" — the failure Harvey reports as the dominant one at
enterprise scale — cannot arise. At 108M it must.

Layout produced under world/corpus/ch/:
    index.sqlite            files(id, matter_id, client_id, folder, filename,
                            ext, bytes, chars, text_path), matters(...)
    text/<matter>/<name>.txt   extracted plain text, written once

Extraction: .docx via python-docx (paragraphs + tables), .xlsx via openpyxl
(sheet grids), .eml via stdlib email (headers + plain body), .pptx via
python-pptx (slide text). Failures are RECORDED in the index with the reason,
never silently dropped — an unparsed file is a fact about the corpus.

Run: python3 world/corpus/build-corpus-index.py --src <dms dir> [--limit N]
"""
from __future__ import annotations

import argparse
import os
import sqlite3
import sys
import time
from email import policy
from email.parser import BytesParser
from pathlib import Path

HERE = Path(__file__).resolve().parent
DEST = HERE / "ch"


def from_docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    out = [x.text for x in d.paragraphs if x.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


def from_xlsx(p: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"[SHEET: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            if not row:
                continue
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append(" | ".join(cells).rstrip(" |"))
    wb.close()
    return "\n".join(out)


def from_eml(p: Path) -> str:
    msg = BytesParser(policy=policy.default).parse(p.open("rb"))
    hdr = [f"{k}: {msg.get(k, '')}" for k in ("From", "To", "Cc", "Date", "Subject")]
    try:
        b = msg.get_body(preferencelist=("plain", "html"))
        body = b.get_content() if b else ""
    except Exception:
        body = ""
    return "\n".join(h for h in hdr if h.split(": ", 1)[1]) + "\n\n" + body


def from_pptx(p: Path) -> str:
    from pptx import Presentation
    pr = Presentation(str(p))
    out = []
    for i, slide in enumerate(pr.slides, 1):
        out.append(f"[SLIDE {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return "\n".join(out)


EXTRACT = {".docx": from_docx, ".xlsx": from_xlsx, ".eml": from_eml, ".pptx": from_pptx,
           ".txt": lambda p: p.read_text(errors="replace"),
           ".md": lambda p: p.read_text(errors="replace")}


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", required=True, help="the dms/ directory to ingest")
    ap.add_argument("--limit", type=int, default=0)
    a = ap.parse_args()

    src = Path(a.src).resolve()
    matters_root = src / "matters" if (src / "matters").is_dir() else src
    DEST.mkdir(parents=True, exist_ok=True)
    (DEST / "text").mkdir(exist_ok=True)

    db = DEST / "index.sqlite"
    if db.exists():
        db.unlink()
    conn = sqlite3.connect(db)
    conn.executescript("""
        CREATE TABLE files (
          id INTEGER PRIMARY KEY, matter_id TEXT, client_id TEXT, folder TEXT,
          filename TEXT, ext TEXT, bytes INTEGER, chars INTEGER,
          text_path TEXT, parse_error TEXT);
        CREATE TABLE matters (
          matter_id TEXT PRIMARY KEY, client_id TEXT, files INTEGER, chars INTEGER);
        CREATE INDEX idx_files_matter ON files(matter_id);
        CREATE INDEX idx_files_ext ON files(ext);
    """)

    paths = sorted(p for p in matters_root.rglob("*") if p.is_file())
    if a.limit:
        paths = paths[:a.limit]
    total = len(paths)
    print(f"ingesting {total} files from {matters_root}", flush=True)

    t0 = time.time()
    ok = err = 0
    per_matter: dict[str, list[int]] = {}
    for i, p in enumerate(paths, 1):
        rel = p.relative_to(matters_root)
        matter_id = rel.parts[0] if len(rel.parts) > 1 else "(root)"
        client_id = matter_id.split("-")[0] if "-" in matter_id else ""
        folder = "/".join(rel.parts[1:-1])
        ext = p.suffix.lower()
        fn = EXTRACT.get(ext)
        text, perr = "", None
        if fn is None:
            perr = f"unsupported ext {ext}"
        else:
            try:
                text = fn(p)
            except Exception as e:
                perr = f"{type(e).__name__}: {e}"[:180]
        tp = ""
        if text:
            outdir = DEST / "text" / matter_id
            outdir.mkdir(parents=True, exist_ok=True)
            safe = rel.name.replace("/", "_") + ".txt"
            (outdir / safe).write_text(text, errors="replace")
            tp = str(Path("text") / matter_id / safe)
        conn.execute(
            "INSERT INTO files (matter_id, client_id, folder, filename, ext, bytes, chars,"
            " text_path, parse_error) VALUES (?,?,?,?,?,?,?,?,?)",
            (matter_id, client_id, folder, rel.name, ext, p.stat().st_size,
             len(text), tp, perr))
        if perr:
            err += 1
        else:
            ok += 1
        per_matter.setdefault(matter_id, [0, 0])
        per_matter[matter_id][0] += 1
        per_matter[matter_id][1] += len(text)
        if i % 500 == 0:
            conn.commit()
            rate = i / max(1e-9, time.time() - t0)
            print(f"  {i}/{total}  ({rate:.0f}/s, {ok} ok, {err} err)", flush=True)

    for m, (n, ch) in per_matter.items():
        conn.execute("INSERT INTO matters (matter_id, client_id, files, chars) VALUES (?,?,?,?)",
                     (m, m.split("-")[0] if "-" in m else "", n, ch))
    conn.commit()

    chars = conn.execute("SELECT COALESCE(SUM(chars),0) FROM files").fetchone()[0]
    print()
    print(f"files ingested : {ok} ok, {err} unparsed")
    print(f"matters        : {len(per_matter)}")
    print(f"clients        : {len({m.split('-')[0] for m in per_matter if '-' in m})}")
    print(f"extracted text : {chars:,} chars  ~{chars // 4:,} tokens")
    if err:
        rows = conn.execute(
            "SELECT ext, COUNT(*) FROM files WHERE parse_error IS NOT NULL GROUP BY ext").fetchall()
        print(f"unparsed by ext: {dict(rows)}")
    conn.close()
    print(f"index          : {db}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
